#!/usr/bin/env python3
from pathlib import Path
from datetime import date
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.section import WD_SECTION
from docx.enum.table import WD_TABLE_ALIGNMENT, WD_CELL_VERTICAL_ALIGNMENT
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from pptx import Presentation
from pptx.util import Inches as PI, Pt as PP
from pptx.dml.color import RGBColor as PRGB
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw, ImageFont
import textwrap

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "courseware"
LABS = ROOT / "labs"
OUT.mkdir(exist_ok=True); LABS.mkdir(exist_ok=True)
TITLE = "AI Vibe Coding for Full Stack Web Development"
CODE = "C1800"
BLUE, TEAL, INK, GREY, PALE = "1F6FEB", "10B981", "161B26", "5B6372", "EAF2FF"

topics = [
 ("1", "Vibe Coding Principles for Full Stack Development", [
  ("Lab 1", "Create an agent-ready project scaffold", "Turn a product brief into a bounded, inspectable full-stack repository.", [
   "Create a TaskFlow repository with frontend, backend, shared, tests and docs folders.",
   "Write PROJECT.md with users, problem, scope, non-goals and acceptance criteria.",
   "Write AGENTS.md with commands, coding conventions, safety constraints and definition of done.",
   "Ask an AI coding agent to propose a file plan before writing code; approve only in-scope files.",
   "Create a vertical-slice backlog: health endpoint, task list, create task and status update.",
   "Run a repository inventory and record the baseline in docs/checkpoints.md."],
   "The agent can explain the architecture and no generated file sits outside the agreed plan."),
  ("Lab 2", "Engineer context and prompts for a vertical slice", "Use goal, context, constraints and verification to generate a stable first slice.", [
   "Create prompts/vertical-slice.md with Goal, Context, Constraints, Deliverables and Verification headings.",
   "Specify a React client, Node/Express API and SQLite persistence; require TypeScript throughout.",
   "Add explicit non-goals: authentication, billing and real-time updates.",
   "Ask the agent to identify assumptions and risks before implementation.",
   "Generate the health endpoint and a frontend connectivity indicator.",
   "Run lint, typecheck and tests; paste results into docs/checkpoints.md.",
   "Start a fresh agent session using only repository context and verify it can continue accurately."],
   "GET /api/health returns 200 and the UI shows API Connected without hidden manual fixes.")]),
 ("2", "Frontend & Backend Implementation with Agents", [
  ("Lab 3", "Generate and refine the frontend", "Build an accessible responsive task dashboard through small agent-reviewed increments.", [
   "Prompt for a component map before code: AppShell, TaskForm, TaskList, TaskCard and StatusFilter.",
   "Generate semantic HTML and keyboard-accessible controls; reject div-only interaction patterns.",
   "Create typed Task and TaskStatus contracts in shared/types.ts.",
   "Implement loading, empty, populated and error states using fixture data.",
   "Add responsive CSS for 375 px and 1280 px widths.",
   "Ask a review agent to inspect accessibility, state coverage and component boundaries.",
   "Run frontend tests and manually check tab order and visible focus."],
   "Users can add and filter fixture tasks at mobile and desktop widths using keyboard controls."),
  ("Lab 4", "Build the API, database and integration tests", "Implement validated routes and persistence behind the typed interface.", [
   "Design the tasks table: id, title, description, status, created_at and updated_at.",
   "Generate a migration and deterministic seed script; inspect SQL before execution.",
   "Implement GET /api/tasks, POST /api/tasks and PATCH /api/tasks/:id/status.",
   "Add schema validation and consistent {error:{code,message,details}} responses.",
   "Keep route, service and repository responsibilities separate.",
   "Connect the frontend through a typed apiClient module and remove fixtures.",
   "Generate integration tests for success, invalid input and missing record cases.",
   "Run the complete test suite twice: clean database and seeded database."],
   "The browser creates and updates persisted tasks; all integration tests pass independently.")]),
 ("3", "Deployment to the Cloud & CI/CD Pipelines", [
  ("Lab 5", "Prepare a production build", "Make configuration, logs, health checks and builds deployment-ready.", [
   "Create .env.example with names and safe placeholders only; keep .env ignored.",
   "Validate required environment variables at application startup.",
   "Generate production build commands for client and server and document their outputs.",
   "Serve the built client and confirm API routing works in production mode.",
   "Add structured request logs without tokens, passwords or personal data.",
   "Add readiness and liveness checks and a graceful shutdown handler.",
   "Run a clean install and production smoke test from a new temporary directory."],
   "A clean checkout builds and starts with documented commands and passes the smoke test."),
  ("Lab 6", "Create CI/CD and release gates", "Automate install, lint, typecheck, test and build on GitHub Actions.", [
   "Prompt the agent for a pipeline threat and failure analysis before YAML generation.",
   "Create .github/workflows/ci.yml using least-privilege permissions and pinned major action versions.",
   "Add dependency caching keyed by the lockfile.",
   "Create separate lint, typecheck, test and build steps with readable failure output.",
   "Upload test reports or build artifacts only when useful; set retention limits.",
   "Protect deployment behind successful CI and environment-specific secrets.",
   "Validate YAML locally, push a feature branch and inspect the Actions run.",
   "Introduce a controlled failing test, observe the gate, then revert it."],
   "The workflow blocks a deliberate defect and passes after the defect is corrected." )]),
 ("4", "Agent Skills, Subagents & MCP Implementation", [
  ("Lab 7", "Create reusable skills and bounded subagents", "Package repeatable routing, security and test workflows without overlapping WSQ tooling.", [
   "Create non-wsq-route-generator with a clear trigger description and minimal instructions.",
   "Create non-wsq-security-review that reports file, line, severity and remediation.",
   "Define non-wsq-frontend-reviewer and non-wsq-api-tester agent roles with read/write boundaries.",
   "Create non-wsq-courseware-qa command and pre/post hooks that reject prohibited programme terms.",
   "Run the skill validator and correct metadata or naming errors.",
   "Delegate frontend review and API test generation as independent bounded tasks.",
   "Review diffs, reconcile conflicts centrally and record decisions."],
   "Every custom artifact begins non-wsq-, passes validation and does not alter WSQ files."),
  ("Lab 8", "Implement an MCP-style tool and orchestrate the capstone", "Expose a safe project-status tool and coordinate verified parallel work.", [
   "Define a project_status tool contract with input schema, output schema and explicit errors.",
   "Implement a local stdio server that reads test summaries and git status without executing arbitrary input.",
   "Register the server in a project-scoped configuration using environment-variable placeholders.",
   "Call the tool with valid input, missing input and malformed input; record structured results.",
   "Ask separate agents to review UI, API, CI and security using bounded prompts.",
   "Merge only evidence-backed changes and run the full verification pipeline.",
   "Generate a release note containing features, test evidence, limitations and rollback steps.",
   "Demonstrate the deployed TaskFlow flow: create, filter and update a task."],
   "The orchestrated release passes all gates and the tool cannot access paths outside the project.")])]

def font(name="Arial", size=11, bold=False, color=INK):
    return {"name":name,"size":Pt(size),"bold":bold,"color":RGBColor.from_string(color)}

def shade(cell, fill):
    tcPr=cell._tc.get_or_add_tcPr(); shd=OxmlElement('w:shd'); shd.set(qn('w:fill'),fill); tcPr.append(shd)

def set_cell(cell, text, bold=False, fill=None):
    cell.text=text; cell.vertical_alignment=WD_CELL_VERTICAL_ALIGNMENT.CENTER
    if fill: shade(cell,fill)
    for p in cell.paragraphs:
        for r in p.runs: r.font.name="Arial"; r.font.size=Pt(9); r.bold=bold

def add_field(paragraph, instruction):
    run=paragraph.add_run(); begin=OxmlElement('w:fldChar'); begin.set(qn('w:fldCharType'),'begin')
    instr=OxmlElement('w:instrText'); instr.set(qn('xml:space'),'preserve'); instr.text=instruction
    sep=OxmlElement('w:fldChar'); sep.set(qn('w:fldCharType'),'separate'); end=OxmlElement('w:fldChar'); end.set(qn('w:fldCharType'),'end')
    run._r.extend([begin,instr,sep,end])

def base_doc(subtitle):
    d=Document(); sec=d.sections[0]; sec.top_margin=Inches(.65); sec.bottom_margin=Inches(.65)
    normal=d.styles['Normal']; normal.font.name='Arial'; normal.font.size=Pt(10.5)
    for s in ['Title','Subtitle','Heading 1','Heading 2','Heading 3']:
        d.styles[s].font.name='Arial'; d.styles[s].font.color.rgb=RGBColor.from_string(BLUE)
    p=d.add_paragraph(); p.alignment=WD_ALIGN_PARAGRAPH.CENTER; p.paragraph_format.space_before=Pt(100)
    r=p.add_run(TITLE); r.bold=True; r.font.name='Arial'; r.font.size=Pt(28); r.font.color.rgb=RGBColor.from_string(BLUE)
    p=d.add_paragraph(); p.alignment=WD_ALIGN_PARAGRAPH.CENTER; r=p.add_run(subtitle); r.bold=True; r.font.size=Pt(19); r.font.color.rgb=RGBColor.from_string(TEAL)
    p=d.add_paragraph(); p.alignment=WD_ALIGN_PARAGRAPH.CENTER; p.add_run(f"Course Code: {CODE}\nVersion 1.0\nConducted by Tertiary Infotech Academy Pte Ltd\nUEN: 201200696W")
    d.add_page_break(); d.add_heading('Document Version Control Record',0)
    t=d.add_table(rows=2, cols=4); t.alignment=WD_TABLE_ALIGNMENT.CENTER; t.style='Table Grid'
    for i,x in enumerate(['Version Number','Effective Date','Summary of Changes','Author']): set_cell(t.rows[0].cells[i],x,True,BLUE)
    for i,x in enumerate(['1.0',str(date.today()),'Initial non-WSQ courseware release','Tertiary Infotech Academy']): set_cell(t.rows[1].cells[i],x)
    d.add_heading('Table of Contents',0); p=d.add_paragraph(); add_field(p,'TOC \\o "1-3" \\h \\z \\u'); d.add_page_break()
    for sec in d.sections:
        f=sec.footer.paragraphs[0]; f.alignment=WD_ALIGN_PARAGRAPH.CENTER
        f.add_run('© 2026 Tertiary Infotech Academy Pte Ltd. All rights reserved.  |  www.tertiarycourses.com.sg  |  Page '); add_field(f,'PAGE'); f.add_run(' of '); add_field(f,'NUMPAGES')
    settings=d.settings._element; update=OxmlElement('w:updateFields'); update.set(qn('w:val'),'true'); settings.append(update)
    return d

def add_note(d,text):
    t=d.add_table(rows=1,cols=1); t.style='Table Grid'; shade(t.cell(0,0),'E8F8F1'); p=t.cell(0,0).paragraphs[0]; p.add_run('Checkpoint: ').bold=True; p.add_run(text)

def build_guide():
    d=base_doc('Learner Guide and Step-by-Step Lab Manual')
    d.add_heading('1. Course Overview',0)
    d.add_paragraph('This two-day, 15-hour practical course teaches an evidence-led workflow for using AI coding agents to architect, implement, test and ship a full-stack application. Learners build one capstone, TaskFlow, through eight connected labs.')
    d.add_heading('Learning outcomes',1)
    for x in ['Frame bounded agent goals and durable project context.','Generate and review accessible frontend components and typed backend services.','Create database models, validation and integration tests.','Build production configuration and CI/CD quality gates.','Create isolated non-WSQ skills, agents, commands and hooks.','Implement and test an MCP-style project tool and orchestrate a release.']:
        d.add_paragraph(x,style='List Bullet')
    d.add_heading('2. Before You Start',0)
    for x in ['Install Git, Node.js 20+, a code editor and one supported coding agent (Codex, Claude Code or Antigravity).','Create a GitHub account and confirm git --version and node --version work.','Create a working folder with no spaces in its path when possible.','Never paste production secrets into prompts. Use .env locally, .env.example for names only, and repository secrets for CI.']:
        d.add_paragraph(x,style='List Number')
    d.add_heading('Prompt contract used in every lab',1)
    d.add_paragraph('Goal — one observable outcome. Context — relevant files and current behaviour. Constraints — boundaries, standards and non-goals. Deliverables — exact files or changes. Verification — commands and acceptance evidence.')
    for num,topic,labs in topics:
        d.add_heading(f'{int(num)+2}. Topic {num}: {topic}',0)
        for lab,name,goal,steps,test in labs:
            d.add_heading(f'{lab}: {name}',1); d.add_paragraph('Goal',style='Heading 2'); d.add_paragraph(goal)
            d.add_paragraph('What you will build',style='Heading 2'); d.add_paragraph('A verified increment of the TaskFlow React + TypeScript + Express + SQLite capstone, plus review evidence in the repository.')
            d.add_paragraph('Procedure',style='Heading 2')
            for i,s in enumerate(steps,1):
                p=d.add_paragraph(); p.paragraph_format.left_indent=Inches(.15); p.paragraph_format.first_line_indent=Inches(-.15); p.add_run(f'{i}.  ').bold=True; p.add_run(s)
                p=d.add_paragraph(); p.paragraph_format.left_indent=Inches(.3); p.add_run('Evidence: ').bold=True; p.add_run('Save the prompt, inspect the diff, run the relevant command, and record the result before continuing.')
            add_note(d,test)
            d.add_paragraph('Reflection',style='Heading 2'); d.add_paragraph('What did the agent assume? Which evidence changed your confidence? What context should be preserved for the next session?')
    d.add_heading('11. Troubleshooting',0)
    for a,b in [('Agent changes too much','Stop, restore the last verified state, narrow deliverables to named files and require a plan first.'),('Types drift across layers','Keep contracts in shared/types.ts and import them from client and server.'),('Tests depend on order','Reset the test database per test and use deterministic seed data.'),('CI passes locally only','Compare Node versions, lockfiles, environment variables and case-sensitive paths.'),('Tool cannot start','Run it directly over stdio, inspect structured error output and verify project-scoped registration.')]:
        p=d.add_paragraph(); p.add_run(a+': ').bold=True; p.add_run(b)
    d.add_heading('12. Glossary',0)
    t=d.add_table(rows=1,cols=2); t.style='Table Grid'; set_cell(t.cell(0,0),'Term',True,BLUE); set_cell(t.cell(0,1),'Meaning',True,BLUE)
    for a,b in [('Vibe coding','Intent-led software development in which a human directs and verifies AI-generated changes.'),('Agent','A model-driven tool that can inspect context, plan and use tools to complete a bounded goal.'),('Context engineering','Selecting durable instructions and task-relevant evidence for reliable agent work.'),('Vertical slice','A small end-to-end feature spanning user interface, API and data.'),('Release gate','An automated check that must pass before deployment.'),('Skill','Reusable instructions and resources triggered for a specialised workflow.'),('Subagent','A separately bounded agent task coordinated by a primary agent.'),('MCP','Model Context Protocol, a standard for exposing tools and context to compatible clients.')]:
        cells=t.add_row().cells; set_cell(cells[0],a,True); set_cell(cells[1],b)
    path=OUT/f'{CODE} {TITLE} - Learner Guide.docx'; d.save(path)
    # aligned learner-facing markdown
    md=[f'# {TITLE}\n\n**Course Code:** {CODE}\n\n**Document:** Learner Guide and Step-by-Step Lab Manual\n', '## Course overview\n\nBuild the TaskFlow capstone through eight verified labs. Use the prompt contract: Goal, Context, Constraints, Deliverables, Verification.\n']
    for num,topic,labs in topics:
        md.append(f'## Topic {num}: {topic}\n')
        for lab,name,goal,steps,test in labs:
            md += [f'### {lab}: {name}\n',f'**Goal:** {goal}\n','**Steps:**\n']+[f'{i}. {s}\n' for i,s in enumerate(steps,1)]+[f'\n**Test it:** {test}\n']
    (ROOT/'LEARNER-GUIDE.md').write_text('\n'.join(md),encoding='utf-8')

def build_labs():
    for num,topic,labs in topics:
        td=LABS/f'topic-{num}'; td.mkdir(exist_ok=True)
        for idx,(lab,name,goal,steps,test) in enumerate(labs,1):
            slug=f'lab-{(int(num)-1)*2+idx}-{name.lower().replace(" ","-").replace("/","-")}'
            content=[f'# {lab}: {name}\n',f'## Goal\n\n{goal}\n','## Prerequisites\n\n- Completed previous lab\n- Development environment verified\n- No real secrets in prompts or repository\n','## Steps\n']
            content += [f'{i}. {s}\n' for i,s in enumerate(steps,1)]
            content += [f'\n## Verification\n\n{test}\n','## Evidence to submit\n\n- Saved prompt\n- Reviewed diff\n- Command/test output\n- Short reflection on one corrected agent assumption\n']
            (td/f'{slug}.md').write_text('\n'.join(content),encoding='utf-8')

def build_lp():
    d=base_doc('Lesson Plan')
    d.add_heading('1. Course Overview',0); d.add_paragraph('Intermediate | 2 days | 15 contact hours | Hands-on instructor-led training | No formal assessment')
    d.add_heading('2. Learning Outcomes',0)
    for x in ['Engineer reliable agent context and prompts.','Build and test a typed full-stack vertical slice.','Automate production checks and delivery gates.','Create reusable non-WSQ agent tooling and an MCP-style integration.']: d.add_paragraph(x,style='List Bullet')
    schedules={
    'Day 1 — Principles and implementation':[
      ('09:30–10:00','Welcome, environment check, goals and responsible AI use','30'),('10:00–11:00','Topic 1: vibe coding, agent tools and architecture','60'),('11:00–12:15','Lab 1: agent-ready scaffold','75'),('12:15–13:00','Lab 2: context and prompt contract, part 1','45'),('13:00–13:30','Lunch','30'),('13:30–14:30','Lab 2: vertical slice and verification','60'),('14:30–15:15','Topic 2: frontend/backend agent workflows','45'),('15:15–15:30','Tea break','15'),('15:30–16:30','Lab 3: accessible responsive frontend','60'),('16:30–17:20','Lab 4: API, data and integration tests','50'),('17:20–17:30','Review and checkpoint','10')],
    'Day 2 — Delivery and orchestration':[
      ('09:30–10:00','Recap, troubleshoot and restore checkpoint','30'),('10:00–10:45','Topic 3: cloud deployment and release gates','45'),('10:45–11:45','Lab 5: production build','60'),('11:45–13:00','Lab 6: GitHub Actions CI/CD','75'),('13:00–13:30','Lunch','30'),('13:30–14:15','CI failure exercise and release review','45'),('14:15–15:00','Topic 4: skills, subagents and MCP','45'),('15:00–15:15','Tea break','15'),('15:15–16:15','Lab 7: non-WSQ reusable agent tooling','60'),('16:15–17:15','Lab 8: MCP-style tool and orchestration','60'),('17:15–17:30','Capstone demonstration and close','15')]}
    d.add_heading('3. Detailed Schedule',0)
    for day,rows in schedules.items():
        d.add_heading(day,1); t=d.add_table(rows=1,cols=3); t.style='Table Grid'; t.alignment=WD_TABLE_ALIGNMENT.CENTER
        for i,x in enumerate(['Time','Topic / Activity','Minutes']): set_cell(t.rows[0].cells[i],x,True,BLUE)
        total=0; instructional=0
        for tm,act,m in rows:
            cs=t.add_row().cells; total += int(m)
            if act != 'Lunch': instructional += int(m)
            fill='E8F8F1' if 'Lab' in act else ('FFF4CC' if 'break' in act.lower() or act=='Lunch' else None)
            for i,x in enumerate([tm,act,m]): set_cell(cs[i],x,False,fill)
        d.add_paragraph(f'Classroom window: {total} minutes; instructional/contact time excluding lunch: {instructional} minutes (7.5 hours).')
        assert total==480 and instructional==450
    d.add_heading('4. Facilitation Methods',0); d.add_paragraph('Short demonstrations, think-aloud prompt design, paired diff reviews, guided labs, failure injection, evidence checkpoints and capstone demonstration.')
    d.add_heading('5. Tools and Resources',0); d.add_paragraph('Laptop, Git, GitHub, Node.js 20+, browser, code editor, supported coding agent, course repository and learner guide. Use synthetic data and placeholder credentials only.')
    d.add_heading('6. Progress Checks',0); d.add_paragraph('Progress is formative: observable lab checkpoints, test output, reviewed diffs and final capstone demonstration. No formal assessment is included.')
    d.save(OUT/f'Lesson Plan - {CODE} {TITLE}.docx')

def add_text(slide,x,y,w,h,text,size=24,bold=False,color=INK,align=PP_ALIGN.LEFT):
    box=slide.shapes.add_textbox(PI(x),PI(y),PI(w),PI(h)); tf=box.text_frame; tf.clear(); tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=align; r=p.add_run(); r.text=text; r.font.name='Arial'; r.font.size=PP(size); r.font.bold=bold; r.font.color.rgb=PRGB.from_string(color); return box

def build_ppt():
    prs=Presentation(); prs.slide_width=PI(13.333); prs.slide_height=PI(7.5)
    def slide(title,kicker=None,body=None):
        s=prs.slides.add_slide(prs.slide_layouts[6]); bg=s.background.fill; bg.solid(); bg.fore_color.rgb=PRGB(255,255,255)
        s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,PI(.14),prs.slide_height).fill.solid(); s.shapes[-1].fill.fore_color.rgb=PRGB.from_string(BLUE); s.shapes[-1].line.fill.background()
        if kicker: add_text(s,.65,.45,11.8,.35,kicker.upper(),11,True,TEAL)
        add_text(s,.65,.85,12,1,title,26,True,INK)
        if body:
            box=add_text(s,.8,2.0,11.7,4.5,'',19,False,GREY); tf=box.text_frame; tf.clear()
            for i,b in enumerate(body):
                p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.text=b; p.level=0; p.font.name='Arial'; p.font.size=PP(18); p.font.color.rgb=PRGB.from_string(GREY); p.space_after=PP(10); p.text='• '+p.text
        add_text(s,.65,7.12,11.5,.2,f'{TITLE} · {CODE}     © 2026 Tertiary Infotech Academy Pte Ltd',9,False,GREY)
        add_text(s,12.1,7.05,.5,.25,str(len(prs.slides)),9,False,GREY,PP_ALIGN.RIGHT); return s
    s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=PRGB(255,255,255)
    add_text(s,.75,.65,3,.4,'TERTIARY INFOTECH ACADEMY',12,True,BLUE); add_text(s,.75,1.7,11.5,1.6,TITLE,34,True,INK); add_text(s,.75,3.5,6,.5,f'Course Code: {CODE}',20,True,TEAL); add_text(s,.75,4.25,8,.8,'Build · Verify · Ship with AI coding agents',24,False,GREY); add_text(s,.75,6.45,9,.4,'2 days · 15 hours · Intermediate · Version 1.0',13,False,GREY)
    slide('Learning outcomes','Course overview',['Engineer goals, prompts and durable context for long agent sessions','Build a typed React, Express and SQLite application','Automate tests, builds and delivery gates','Package non-WSQ skills, agents, commands and hooks','Implement an MCP-style tool and orchestrate a verified release'])
    slide('The capstone: TaskFlow','Course overview',['One application grows through all eight labs','Frontend: accessible responsive React + TypeScript','Backend: Express API with validation','Data: SQLite migrations and deterministic seed','Delivery: GitHub Actions and production checks'])
    slide('The prompt contract','Core practice',['Goal — one observable outcome','Context — only relevant files and current behaviour','Constraints — standards, boundaries and non-goals','Deliverables — explicit files and changes','Verification — commands and acceptance evidence'])
    for num,topic,labs in topics:
        slide(f'Topic {num}',topic,[f'{lab}: {name}' for lab,name,_,_,_ in labs])
        concepts={"1":['Vibe coding is intent-led, but evidence-controlled','Agents differ from autocomplete: they plan, inspect and use tools','Repository context outlives chat context','Prefer vertical slices and small reviewed diffs'],"2":['Shared contracts prevent frontend/backend drift','Generate states, not only happy-path screens','Separate route, service and repository layers','Tests specify behaviour and constrain regeneration'],"3":['Configuration is code; secrets are not','A clean build is stronger evidence than “works on my machine”','CI gates should be deterministic and least-privileged','Deployment needs health checks, logs and rollback'],"4":['Skills encode repeatable workflow knowledge','Subagents need bounded tasks and central reconciliation','MCP exposes typed tools and context','Project-scoped namespacing prevents cross-programme interference']}[num]
        slide('Key concepts',f'Topic {num}',concepts)
        for lab,name,goal,steps,test in labs:
            slide(f'{lab}: {name}',f'Topic {num}',[goal,'Deliverable: a verified TaskFlow increment','Evidence: prompt + reviewed diff + command output'])
            for i in range(0,len(steps),3): slide(f'{lab} · Steps {i+1}–{min(i+3,len(steps))}','Guided lab',[f'{j+1}. {steps[j]}' for j in range(i,min(i+3,len(steps)))])
            slide('Test it',lab,[test,'If it fails: stop, narrow the task, inspect evidence, then retry.'])
    slide('Capstone release checklist','Course close',['All lint, typecheck, unit and integration checks pass','Production build starts from a clean checkout','CI blocks a deliberate defect','No secrets or prohibited programme content','Release note states limitations and rollback steps','TaskFlow create, filter and update flow succeeds'])
    slide('Keep the human in the loop','Course close',['Direct with explicit intent','Inspect plans and diffs','Verify with repeatable evidence','Preserve useful context','Ship small, reversible changes'])
    prs.save(OUT/f'{CODE} {TITLE}.pptx')

def screenshot():
    im=Image.new('RGB',(1600,900),'white'); d=ImageDraw.Draw(im)
    try: big=ImageFont.truetype('/System/Library/Fonts/Supplemental/Arial Bold.ttf',64); med=ImageFont.truetype('/System/Library/Fonts/Supplemental/Arial.ttf',30)
    except: big=med=None
    d.rectangle((0,0,28,900),fill='#1F6FEB'); d.text((100,110),'C1800',fill='#10B981',font=med); d.multiline_text((100,190),'AI Vibe Coding for\nFull Stack Web Development',fill='#161B26',font=big,spacing=18)
    d.text((100,465),'Complete non-WSQ courseware · 8 connected labs · 2 days',fill='#5B6372',font=med)
    for i,t in enumerate(['CONTEXT','FULL STACK','CI/CD','SKILLS + MCP']):
        x=100+i*360; d.rounded_rectangle((x,610,x+315,735),20,fill='#EAF2FF',outline='#1F6FEB',width=3); d.text((x+28,655),t,fill='#1F6FEB',font=med)
    d.text((100,820),'Tertiary Infotech Academy Pte Ltd',fill='#5B6372',font=med); im.save(ROOT/'screenshot.png')

if __name__=='__main__':
    build_labs(); build_guide(); build_lp(); build_ppt(); screenshot(); print('Courseware generated')
