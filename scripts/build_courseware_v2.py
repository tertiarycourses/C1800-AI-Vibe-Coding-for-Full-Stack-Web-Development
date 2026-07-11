#!/usr/bin/env python3
"""High-depth C1800 non-WSQ courseware generator."""
from pathlib import Path
from datetime import date
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from pptx import Presentation
from pptx.util import Inches as I, Pt as P
from pptx.dml.color import RGBColor as C
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import re, shutil

ROOT=Path(__file__).resolve().parents[1]; OUT=ROOT/'courseware'; LABROOT=ROOT/'labs'
TITLE='AI Vibe Coding for Full Stack Web Development'; CODE='C1800'
BLUE='1F6FEB'; TEAL='10B981'; INK='161B26'; GREY='5B6372'; LIGHT='F5F8FC'; VIOLET='7C3AED'
REPO='https://github.com/tertiarycourses/C1800-AI-Vibe-Coding-for-Full-Stack-Web-Development'
LOGO=OUT/'assets/tertiary-infotech-logo.png'
VISUALS=OUT/'assets/visuals'

topics=[
('01','Vibe Coding Principles for Full Stack Development',
 ['Vibe coding as an engineering workflow','Agentic tools versus autocomplete','Goal and context engineering','Full-stack architecture and vertical slices','Long-session memory and checkpoints'],[
 ('1.1','Establish the Agentic Engineering Loop','Create a visible seven-stage loop and use it to correct a deliberately vague request.','docs/agentic-loop.md','Write the loop, improve a weak prompt, capture assumptions, and define evidence before code.'),
 ('1.2','Set Up and Verify the Development Workbench','Install and verify Node.js, Git, an editor, browser tools and one coding agent.','docs/environment-check.md','Record versions, run a diagnostic project, and resolve one predictable setup failure.'),
 ('1.3','Compare Codex, Claude Code and Antigravity','Run the same bounded repository-analysis task in available tools and compare evidence.','docs/tool-comparison.md','Compare plan quality, file awareness, tool use, diff clarity and verification discipline.'),
 ('1.4','Turn a Business Need into a Product Brief','Convert a TaskFlow scenario into users, jobs, scope, non-goals and acceptance criteria.','PROJECT.md','Create a decision-ready brief that a fresh agent can interpret without chat history.'),
 ('1.5','Create Durable Repository Instructions','Write AGENTS.md and project rules that guide agents without over-constraining implementation.','AGENTS.md','Specify commands, architecture, conventions, safety rules and definition of done.'),
 ('1.6','Architect and Scaffold a Full-Stack Vertical Slice','Plan React, Express, SQLite and shared contracts; scaffold only after reviewing the plan.','apps/web + apps/api + packages/shared','Create a runnable monorepo skeleton and a health-check slice with no hidden manual fixes.')]),
('02','Frontend & Backend Implementation with Agents',
 ['Component and API boundaries','Accessible responsive UI','Typed contracts and validation','Persistence and migrations','Testing as executable context','Debugging agent-generated defects'],[
 ('2.1','Design the TaskFlow UI System','Translate the product brief into tokens, components, states and responsive wireframes.','apps/web/src/styles + docs/ui-system.md','Define accessible colour, spacing, typography and interaction states before component generation.'),
 ('2.2','Generate the Application Shell and Navigation','Build a semantic responsive shell with skip link, header, navigation and main region.','apps/web/src/components/AppShell.tsx','Inspect landmarks, tab order, focus styles and mobile layout.'),
 ('2.3','Build Task Cards, Lists and Empty States','Create typed presentational components for loading, empty, populated and error states.','apps/web/src/components/tasks','Use stable keys, clear status semantics and component-level examples.'),
 ('2.4','Build a Controlled Task Form','Implement accessible controlled inputs, client validation and submission feedback.','apps/web/src/components/TaskForm.tsx','Handle labels, errors, disabled state, focus recovery and double-submit prevention.'),
 ('2.5','Model Shared TypeScript Contracts','Define Task, CreateTaskInput, TaskStatus and API result types once for both layers.','packages/shared/src/task.ts','Remove duplicated shapes and verify client/server compilation against shared contracts.'),
 ('2.6','Create the Express API Structure','Generate app, routes, controllers, services and middleware with explicit boundaries.','apps/api/src','Implement health and task routes without placing business logic in route handlers.'),
 ('2.7','Add Runtime Validation and Safe Errors','Validate untrusted requests and return consistent errors without leaking internals.','apps/api/src/middleware + schemas','Test malformed JSON, blank titles, invalid status and unknown records.'),
 ('2.8','Add SQLite Migrations and Repository Logic','Create schema, migration, seed and repository functions with parameterised queries.','apps/api/src/db','Prove clean migration, deterministic seed, CRUD persistence and SQL-injection resistance.'),
 ('2.9','Connect the React Client to the API','Replace fixtures with a typed API client and resilient loading/error behaviour.','apps/web/src/lib/apiClient.ts','Exercise create, read and update across the browser, API and database.'),
 ('2.10','Generate Unit and Integration Tests','Use tests as executable context for components, services and HTTP endpoints.','apps/web/src/**/*.test.tsx + apps/api/test','Cover happy paths, boundaries and failures; ensure tests are isolated and deterministic.'),
 ('2.11','Debug an Agent-Generated Cross-Layer Defect','Trace a status mismatch from UI symptom through network, validation and persistence.','docs/debug-report.md','Form hypotheses, collect evidence, make the smallest fix and add a regression test.')]),
('03','Deployment to the Cloud & CI/CD Pipelines',
 ['Production configuration','Repeatable builds','Containers and health checks','CI quality gates','Preview and production deployment','Observability and rollback'],[
 ('3.1','Harden Configuration and Secret Handling','Validate environment variables, create safe examples and prevent secret exposure.','.env.example + apps/api/src/config.ts','Fail fast with helpful messages while keeping credentials out of prompts, logs and Git.'),
 ('3.2','Create and Test Production Builds','Build client and server from a clean checkout and run a production smoke test.','dist outputs + docs/production-runbook.md','Document exact commands, artifacts, runtime assumptions and recovery steps.'),
 ('3.3','Containerise the Full-Stack Application','Generate a multi-stage Dockerfile, health check and minimal runtime image.','Dockerfile + .dockerignore','Build, run, inspect, stop and rebuild without copying secrets or development clutter.'),
 ('3.4','Build a GitHub Actions CI Pipeline','Create least-privilege lint, typecheck, test and build gates with caching.','.github/workflows/ci.yml','Validate YAML, observe a passing run and preserve readable failure evidence.'),
 ('3.5','Prove the Release Gate with Failure Injection','Introduce controlled lint, type and test failures and confirm CI blocks each one.','docs/ci-failure-lab.md','Distinguish product defects from pipeline defects and revert cleanly after evidence.'),
 ('3.6','Deploy a Preview and Production Release','Configure a cloud target, environment variables, migrations and gated deployment.','deployment config + release URL','Verify health, UI/API connectivity, persistence and environment isolation.'),
 ('3.7','Add Observability, Release Notes and Rollback','Add structured logs, correlation IDs, health endpoints and a rollback runbook.','docs/releases + docs/rollback.md','Detect a simulated failure, identify the release, and execute a reversible rollback decision.')]),
('04','Agent Skills, Subagents & MCP Implementation',
 ['Reusable skills and trigger design','Commands and lifecycle hooks','Bounded subagent delegation','Model Context Protocol tools','Security boundaries','Orchestrated release engineering'],[
 ('4.1','Create a Non-WSQ Route-Generator Skill','Package a repeatable typed-route workflow with concise triggers and validation.','.codex/skills/non-wsq-route-generator','Create, validate and invoke the skill on one new endpoint.'),
 ('4.2','Create Commands and Quality Hooks','Add project commands plus pre/post hooks that enforce scope and verification.','.claude/commands + .claude/hooks','Reject prohibited programme content and incomplete lab artifacts without touching WSQ tooling.'),
 ('4.3','Design Bounded Review and Test Subagents','Define frontend-review, API-test and security-review roles with explicit outputs.','.claude/agents/non-wsq-*','Delegate independent tasks, compare findings and keep integration decisions central.'),
 ('4.4','Build a Local MCP Project-Status Server','Expose typed read-only tools for Git status, test summaries and release readiness.','tools/project-status-mcp','Implement stdio transport, schemas, structured errors and path containment.'),
 ('4.5','Test MCP Tools and Security Boundaries','Exercise valid, invalid and adversarial inputs without allowing arbitrary execution.','tools/project-status-mcp/test','Prove schema validation, project-root containment, output limits and safe error messages.'),
 ('4.6','Orchestrate the TaskFlow Capstone Release','Coordinate specialised reviews through the agentic loop and ship a verified release.','docs/capstone-release.md','Reconcile findings, run all gates, demonstrate TaskFlow and publish evidence-backed release notes.')])]

def slug(s): return re.sub(r'[^a-z0-9]+','-',s.lower()).strip('-')
def all_labs():
 for tn,topic,concepts,labs in topics:
  for lab in labs: yield tn,topic,lab

def field(p,code):
 r=p.add_run(); a=OxmlElement('w:fldChar'); a.set(qn('w:fldCharType'),'begin'); b=OxmlElement('w:instrText'); b.set(qn('xml:space'),'preserve'); b.text=code; c=OxmlElement('w:fldChar'); c.set(qn('w:fldCharType'),'end'); r._r.extend([a,b,c])
def shade(cell,color):
 pr=cell._tc.get_or_add_tcPr(); e=OxmlElement('w:shd'); e.set(qn('w:fill'),color); pr.append(e)
def base_doc(subtitle):
 d=Document(); sec=d.sections[0]; sec.top_margin=Inches(.7); sec.bottom_margin=Inches(.65)
 for st in d.styles:
  if hasattr(st,'font'): st.font.name='Arial'
 d.styles['Normal'].font.size=Pt(10.5)
 for n,s,c in [('Title',30,BLUE),('Heading 1',20,BLUE),('Heading 2',16,TEAL),('Heading 3',12,VIOLET)]: d.styles[n].font.size=Pt(s); d.styles[n].font.bold=True; d.styles[n].font.color.rgb=RGBColor.from_string(c)
 if LOGO.exists():
  p=d.add_paragraph(); p.alignment=WD_ALIGN_PARAGRAPH.CENTER; p.add_run().add_picture(str(LOGO),width=Inches(1.05))
 p=d.add_paragraph(); p.alignment=WD_ALIGN_PARAGRAPH.CENTER; p.paragraph_format.space_before=Pt(25); r=p.add_run(TITLE); r.bold=True; r.font.size=Pt(29); r.font.color.rgb=RGBColor.from_string(BLUE)
 p=d.add_paragraph(); p.alignment=WD_ALIGN_PARAGRAPH.CENTER; r=p.add_run(subtitle); r.bold=True; r.font.size=Pt(19); r.font.color.rgb=RGBColor.from_string(TEAL)
 p=d.add_paragraph(); p.alignment=WD_ALIGN_PARAGRAPH.CENTER; p.add_run(f'Course Code: {CODE}\nVersion 2.0\nConducted by Tertiary Infotech Academy Pte Ltd\nUEN: 201200696W')
 d.add_page_break(); d.add_heading('Document Version Control Record',0); t=d.add_table(rows=3,cols=4); t.style='Table Grid'; t.alignment=WD_TABLE_ALIGNMENT.CENTER
 for i,x in enumerate(['Version','Effective Date','Summary of Changes','Author']): t.cell(0,i).text=x; shade(t.cell(0,i),BLUE)
 for row,data in enumerate([('1.0','11 July 2026','Initial release','Tertiary Infotech Academy'),('2.0',str(date.today()),'Rebuilt to 24-lab adult-training standard','Tertiary Infotech Academy')],1):
  for i,x in enumerate(data): t.cell(row,i).text=x
 d.add_heading('Table of Contents',0); field(d.add_paragraph(),'TOC \\o "1-3" \\h \\z \\u'); d.add_page_break()
 for sec in d.sections:
  p=sec.footer.paragraphs[0]; p.alignment=WD_ALIGN_PARAGRAPH.CENTER; p.add_run('© 2026 Tertiary Infotech Academy Pte Ltd. All rights reserved. | www.tertiarycourses.com.sg | Page '); field(p,'PAGE'); p.add_run(' of '); field(p,'NUMPAGES')
 u=OxmlElement('w:updateFields'); u.set(qn('w:val'),'true'); d.settings._element.append(u); return d

def para(d,text,boldlead=None):
 p=d.add_paragraph();
 if boldlead and text.startswith(boldlead): p.add_run(boldlead).bold=True; p.add_run(text[len(boldlead):])
 else: p.add_run(text)
 return p
def note(d,label,text,color='E8F8F1'):
 t=d.add_table(rows=1,cols=1); t.style='Table Grid'; shade(t.cell(0,0),color); p=t.cell(0,0).paragraphs[0]; p.add_run(label+': ').bold=True; p.add_run(text)
def code(d,text):
 t=d.add_table(rows=1,cols=1); shade(t.cell(0,0),'EEF2F7'); p=t.cell(0,0).paragraphs[0]; r=p.add_run(text); r.font.name='Courier New'; r.font.size=Pt(8.5)

def lab_steps(labno,name,goal,target,practice):
 return [
  ('Frame',f'Read the scenario and restate the outcome in one sentence: “{goal}” Identify the learner-visible success condition and write it under Goal in `prompts/lab-{labno}.md`.'),
  ('Inspect the starting state',f'Run `git status --short`, inventory the relevant files, and open the current checkpoint. Confirm that the intended change is limited to `{target}`. If unrelated changes exist, record them and do not overwrite them.'),
  ('Plan',f'Ask the coding agent for a numbered implementation plan. Require named files, data flow, risks, and verification commands. Do not permit edits yet. Compare the proposal with the architecture and remove unnecessary changes.'),
  ('Create the prompt contract',f'Write Goal, Context, Constraints, Deliverables and Verification sections. Include the requirement: {practice} Add explicit non-goals so the agent cannot silently expand scope.'),
  ('Generate a small increment',f'Authorise only the first coherent change in `{target}`. Ask the agent to explain each file before editing it. Keep the diff small enough to review in under five minutes.'),
  ('Inspect the diff',f'Run `git diff -- {target.split(" + ")[0]}` where applicable. Check correctness, naming, accessibility/security, error handling, duplicated logic and unexpected dependencies. Reject code you cannot explain.'),
  ('Verify',f'Run the lab-specific command plus `npm run lint`, `npm run typecheck` and the narrowest relevant test. Capture the command, exit status and meaningful output in `docs/checkpoints.md`.'),
  ('Correct',f'If verification fails, give the agent the exact error, observed behaviour and relevant file—not a guess. Ask for two hypotheses, choose one test, and apply the smallest evidence-backed correction.'),
  ('Challenge',f'Change one boundary condition (empty input, network failure, unknown ID, narrow viewport or missing configuration). Predict the behaviour before running it, then add or update a regression test.'),
  ('Commit the checkpoint',f'Confirm `git diff --check` and a clean secret scan. Commit the verified increment with a precise message such as `feat(lab-{labno}): {slug(name)}`. Record one assumption corrected during the lab.')]

def build_labs_and_md():
 if LABROOT.exists(): shutil.rmtree(LABROOT)
 md=[f'# {TITLE} — Learner Guide\n',f'**Course Code:** {CODE}\n\n**Version:** 2.0\n\n**Source repository:** {REPO}\n',
 '## How to use this guide\n\nThis guide is a working manual. Type the commands, save every prompt, inspect every diff, and keep evidence. Do not accept agent output merely because it looks plausible. Each lab follows the adult-learning rhythm: trainer demonstration, guided practice, independent challenge, reflection.\n',
 '## The Agentic AI Loop\n\n**Frame → Plan → Generate → Inspect → Verify → Correct → Commit.** The loop converts natural-language intent into controlled engineering work. Frame makes success observable. Plan exposes assumptions before code. Generate keeps changes small. Inspect applies human judgement. Verify produces repeatable evidence. Correct uses diagnosis instead of prompt roulette. Commit preserves a reversible checkpoint.\n',
 '## Before You Start\n\nInstall Node.js 20 or later, Git, Docker Desktop, a modern browser, VS Code, and at least one supported agentic coding tool. Create a GitHub account and a cloud deployment account. Use only synthetic training data. Never paste keys, tokens, passwords, personal data or production source code into an unapproved AI service.\n',
 '### Environment verification\n\n```bash\nnode --version\nnpm --version\ngit --version\ndocker --version\ngh auth status\n```\n\nExpected result: each command prints a version; GitHub CLI shows an authenticated account. If Docker is unavailable, complete container steps as a trainer demonstration and continue with local Node processes.\n']
 for tn,topic,concepts,labs in topics:
  md += [f'## Topic {tn} — {topic}\n', '### Why this topic matters\n\n'+(' '.join(concepts))+'. An adult practitioner must be able to explain the trade-off, demonstrate the workflow, diagnose predictable failure and transfer it to a new project.\n','### Key concepts\n']
  for c in concepts: md.append(f'- **{c}.** Connect the concept to an observable engineering decision; ask what evidence would prove the decision correct.\n')
  md.append('### Deep dive\n\n')
  md.append(f'{topic} is not a collection of isolated features. It is a decision system. Start with the user outcome, make boundaries explicit, preserve shared contracts, and use automation to reveal errors early. AI agents accelerate drafting and exploration, but responsibility for architecture, privacy, accessibility, security and release readiness remains with the human team.\n\nUse a two-pass review. In the first pass, check whether the change solves the stated problem. In the second, examine failure behaviour: invalid data, unavailable dependencies, retries, concurrency, narrow screens, missing configuration and partial deployment. This separation prevents polished code from distracting you from a wrong solution.\n')
  td=LABROOT/f'topic-{tn}'; td.mkdir(parents=True,exist_ok=True)
  for labno,name,goal,target,practice in labs:
   steps=lab_steps(labno,name,goal,target,practice)
   prompt=f'''You are working on TaskFlow, a training full-stack application.\n\nGOAL\n{goal}\n\nCONTEXT\nRelevant target: {target}\nCurrent checkpoint: verified previous lab\n\nCONSTRAINTS\n- Propose a plan before editing.\n- Keep the change within the named target.\n- Preserve TypeScript strictness, accessibility and safe error handling.\n- Do not add credentials, broad dependencies or unrelated refactors.\n\nDELIVERABLES\n- Implementation for: {practice}\n- Focused tests and documentation updates\n\nVERIFICATION\nList exact commands and expected observable behaviour. Report assumptions and remaining risks.'''
   section=[f'### Lab {labno} — {name}\n',f'**Outcome:** {goal}\n\n**Primary target:** `{target}`\n',
    '**Estimated time:** 25–40 minutes\n\n**Learning mode:** trainer demonstration → guided pair work → independent boundary challenge\n',
    '#### Starting checkpoint\n\nOpen the previous lab checkpoint. Confirm the application starts and the last lab verification passes. If it does not, stop and restore that checkpoint before adding new work.\n',
    '#### Reference prompt\n\n```text\n'+prompt+'\n```\n',
    '#### Step-by-step procedure\n']
   for i,(phase,detail) in enumerate(steps,1):
    section += [f'**Step {i} — {phase}**\n\n{detail}\n',f'> Evidence: save the prompt or command output under `evidence/lab-{labno}/step-{i}`. Write one sentence explaining what the evidence proves.\n']
   section += ['#### Expected result\n\n'+practice+' The repository remains understandable to a fresh agent, relevant checks pass, and no unrelated files change.\n',
    '#### Troubleshooting and recovery\n\n- **The agent edits too many files:** stop the run, keep the evidence, restore only agent-created changes, name the permitted files, and restart at Plan.\n- **A command is missing:** inspect `package.json` before inventing a replacement; add the script only if it is part of the agreed deliverable.\n- **The test passes but the feature fails:** verify that the test reaches the real boundary and is not mocking the behaviour under investigation.\n- **The error changes after every prompt:** stop prompt iteration. Reproduce once, minimise the case, form two hypotheses and collect discriminating evidence.\n',
    '#### Independent challenge\n\nRepeat the verification with a changed boundary condition. Do not ask the agent for the answer first. Write your prediction, run the experiment, then use the discrepancy to improve the implementation or test.\n',
    '#### Reflection\n\n1. Which assumption had the highest risk? 2. What evidence increased confidence most? 3. What should be preserved in project context? 4. What would you explain differently to a colleague?\n']
   md.extend(section)
   (td/f'lab-{labno}-{slug(name)}.md').write_text('\n'.join(section),encoding='utf-8')
 md += ['## Appendix A — Prompt Review Checklist\n\nBefore sending: Is the goal observable? Is context sufficient but bounded? Are constraints meaningful? Are deliverables named? Are verification commands independent of the agent? Are secrets and personal data excluded?\n',
 '## Appendix B — AI-Generated Code Review Checklist\n\nCheck scope, behaviour, types, error paths, accessibility, security, tests, dependency necessity, configuration, logs, performance, readability and rollback. Trace data from user input to storage and back.\n',
 '## Appendix C — Troubleshooting Decision Tree\n\nReproduce → minimise → classify layer → gather evidence → state two hypotheses → run one discriminating test → apply smallest fix → add regression test → rerun full gate.\n',
 '## Appendix D — Glossary\n\n**Agentic loop:** controlled cycle from intent to verified checkpoint. **Context engineering:** deliberate selection of instructions and evidence. **Vertical slice:** thin end-to-end user value. **Contract:** shared interface between components. **Release gate:** automated condition for promotion. **Skill:** reusable workflow knowledge. **Subagent:** bounded delegated agent task. **MCP:** Model Context Protocol for typed tools and context.\n']
 text='\n'.join(md); (ROOT/'LEARNER-GUIDE.md').write_text(text,encoding='utf-8'); (OUT/'C1800-Learner-Guide.md').write_text(text,encoding='utf-8')
 return md

def build_guide(md):
 d=base_doc('Learner Guide and Detailed Lab Manual')
 in_code=False; buf=[]
 for raw in '\n'.join(md).splitlines():
  if raw.startswith('```'):
   if in_code: code(d,'\n'.join(buf)); buf=[]
   in_code=not in_code; continue
  if in_code: buf.append(raw); continue
  if raw.startswith('# '): d.add_heading(raw[2:],0)
  elif raw.startswith('## '): d.add_heading(raw[3:],0)
  elif raw.startswith('### '): d.add_heading(raw[4:],1)
  elif raw.startswith('#### '): d.add_heading(raw[5:],2)
  elif raw.startswith('**Step '): d.add_heading(raw.replace('**',''),2)
  elif raw.startswith('> '): note(d,'Evidence',raw[2:].replace('Evidence: ','').strip(),'EAF2FF')
  elif raw.startswith('- '): d.add_paragraph(raw[2:],style='List Bullet')
  elif re.match(r'^\d+\. ',raw): d.add_paragraph(raw,style='List Number')
  elif raw.strip(): para(d,re.sub(r'\*\*([^*]+)\*\*',r'\1',raw))
 # deliberate page capacity: each lab begins near a new page through heading pagination
 for p in d.paragraphs:
  if p.text.startswith('Lab ') and ' — ' in p.text: p.paragraph_format.page_break_before=True
 d.save(OUT/f'{CODE} {TITLE} - Learner Guide.docx')

def tb(slide,x,y,w,h,text,size=22,bold=False,color=INK,align=PP_ALIGN.LEFT):
 b=slide.shapes.add_textbox(I(x),I(y),I(w),I(h)); f=b.text_frame; f.clear(); f.word_wrap=True; p=f.paragraphs[0]; p.alignment=align; r=p.add_run(); r.text=text; r.font.name='Arial'; r.font.size=P(size); r.font.bold=bold; r.font.color.rgb=C.from_string(color); return b
def build_ppt():
 prs=Presentation(); prs.slide_width=I(13.333); prs.slide_height=I(7.5)
 palette=[BLUE,TEAL,VIOLET,'F59E0B','EF4444']
 pale=['EAF2FF','E8F8F1','F1EAFE','FFF4CC','FDECEC']
 def shape(s,kind,x,y,w,h,fill,line=None,radius=True):
  q=s.shapes.add_shape(kind,I(x),I(y),I(w),I(h)); q.fill.solid(); q.fill.fore_color.rgb=C.from_string(fill); q.line.color.rgb=C.from_string(line or fill); return q
 def shell(title,kicker='COURSE'):
  s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=C(255,255,255)
  sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,I(.14),prs.slide_height); sh.fill.solid(); sh.fill.fore_color.rgb=C.from_string(BLUE); sh.line.fill.background()
  # subtle visual anchor and progress rail
  shape(s,MSO_SHAPE.OVAL,11.85,-.32,1.15,1.15,'F1F6FF','F1F6FF'); shape(s,MSO_SHAPE.OVAL,12.42,.18,.42,.42,'D9F7EA','D9F7EA')
  tb(s,.65,.35,12,.3,kicker.upper(),11,True,TEAL); tb(s,.65,.78,12,1,title,28,True)
  tb(s,.65,7.12,11.7,.18,f'{TITLE} · {CODE}     © 2026 Tertiary Infotech Academy Pte Ltd',8,False,GREY); tb(s,12.25,7.08,.4,.2,str(len(prs.slides)),8,False,GREY,PP_ALIGN.RIGHT); return s
 def visual_story(title,kicker,image_name,takeaways):
  s=shell(title,kicker); img=VISUALS/image_name
  if img.exists(): s.shapes.add_picture(str(img),I(.7),I(1.68),I(8.15),I(4.58))
  shape(s,MSO_SHAPE.ROUNDED_RECTANGLE,9.08,1.68,3.52,4.58,'F5F8FC','D7E3F5')
  tb(s,9.42,1.98,2.8,.25,'LOOK → NOTICE → EXPLAIN',10,True,TEAL)
  for i,x in enumerate(takeaways):
   y=2.55+i*1.05; shape(s,MSO_SHAPE.OVAL,9.4,y,.48,.48,palette[i]); tb(s,9.4,y+.12,.48,.18,str(i+1),10,True,'FFFFFF',PP_ALIGN.CENTER); tb(s,10.05,y-.02,2.15,.72,x,13,True,INK)
 def bullets(title,kicker,items):
  s=shell(title,kicker); items=items[:5]
  # Prompt contracts and core loops become connected process diagrams.
  if 'prompt contract' in title.lower() or 'agentic ai engineering loop' in title.lower() or 'repeatable practice' in title.lower():
   n=len(items); gap=.16; w=(11.75-gap*(n-1))/n
   for i,x in enumerate(items):
    x0=.72+i*(w+gap); shape(s,MSO_SHAPE.ROUNDED_RECTANGLE,x0,2.05,w,2.95,pale[i],palette[i])
    shape(s,MSO_SHAPE.OVAL,x0+.15,2.28,.52,.52,palette[i]); tb(s,x0+.15,2.4,.52,.2,str(i+1),11,True,'FFFFFF',PP_ALIGN.CENTER)
    parts=x.split(' — ',1); tb(s,x0+.18,2.95,w-.36,.48,parts[0],15,True,palette[i],PP_ALIGN.CENTER)
    if len(parts)>1: tb(s,x0+.18,3.52,w-.36,1.16,parts[1],11,False,GREY,PP_ALIGN.CENTER)
    if i<n-1: shape(s,MSO_SHAPE.CHEVRON,x0+w-.02,3.18,.28,.48,'BFD4F7','BFD4F7')
  # Procedures use numbered action cards instead of dense bullets.
  elif title.lower().startswith('procedure'):
   positions=[(.75,1.85,5.75,2.05),(6.75,1.85,5.75,2.05),(.75,4.08,5.75,2.05),(6.75,4.08,5.75,2.05)]
   for i,x in enumerate(items):
    px,py,pw,ph=positions[i]; shape(s,MSO_SHAPE.ROUNDED_RECTANGLE,px,py,pw,ph,'FFFFFF','D7E3F5')
    shape(s,MSO_SHAPE.RECTANGLE,px,py,.12,ph,palette[i]); m=re.match(r'(\d+)\. ([^—]+)—\s*(.*)',x)
    num,head,body=(m.group(1),m.group(2).strip(),m.group(3).strip()) if m else (str(i+1),'ACTION',x)
    shape(s,MSO_SHAPE.OVAL,px+.25,py+.24,.55,.55,palette[i]); tb(s,px+.25,py+.37,.55,.2,num,11,True,'FFFFFF',PP_ALIGN.CENTER)
    tb(s,px+.95,py+.22,pw-1.2,.38,head,15,True,palette[i]); tb(s,px+.28,py+.94,pw-.56,.88,body,12,False,GREY)
  # Lab overview becomes a goal panel plus a build/evidence stack.
  elif 'lab overview' in kicker.lower():
   shape(s,MSO_SHAPE.ROUNDED_RECTANGLE,.72,1.78,5.15,4.72,'EAF2FF',BLUE); tb(s,1.02,2.05,1.4,.25,'YOUR MISSION',11,True,BLUE); tb(s,1.02,2.55,4.55,2.2,items[0],22,True,INK); tb(s,1.02,5.35,4.4,.35,'BUILD A VERIFIED INCREMENT',12,True,TEAL)
   labels=['TARGET','BUILD','EVIDENCE','TRANSFER']
   for i,x in enumerate(items[1:5]):
    y=1.78+i*1.14; shape(s,MSO_SHAPE.ROUNDED_RECTANGLE,6.15,y,6.35,.96,'FFFFFF','D7E3F5'); shape(s,MSO_SHAPE.RECTANGLE,6.15,y,.11,.96,palette[i]); tb(s,6.45,y+.15,1.05,.22,labels[i],10,True,palette[i]); tb(s,7.48,y+.12,4.7,.54,x,13,False,INK)
  # Verification slides are visual checklists.
  elif 'expected result' in title.lower() or 'release review' in title.lower():
   shape(s,MSO_SHAPE.ROUNDED_RECTANGLE,.72,1.78,3.05,4.75,'E8F8F1',TEAL); shape(s,MSO_SHAPE.OVAL,1.53,2.35,1.4,1.4,TEAL); tb(s,1.53,2.72,1.4,.5,'✓',34,True,'FFFFFF',PP_ALIGN.CENTER); tb(s,1.05,4.18,2.38,.8,'EVIDENCE\nBEATS CONFIDENCE',16,True,TEAL,PP_ALIGN.CENTER)
   for i,x in enumerate(items):
    y=1.82+i*.9; shape(s,MSO_SHAPE.ROUNDED_RECTANGLE,4.05,y,8.43,.7,'FFFFFF','D7E3F5'); shape(s,MSO_SHAPE.OVAL,4.28,y+.15,.38,.38,TEAL); tb(s,4.28,y+.23,.38,.16,'✓',10,True,'FFFFFF',PP_ALIGN.CENTER); tb(s,4.85,y+.13,7.25,.35,x,14,False,INK)
  # Troubleshooting is a diagnostic path.
  elif 'troubleshoot' in title.lower() or 'failure pattern' in title.lower():
   for i,x in enumerate(items):
    y=1.82+i*.9; shape(s,MSO_SHAPE.ROUNDED_RECTANGLE,1.15,y,10.85,.68,pale[i],palette[i]); shape(s,MSO_SHAPE.OVAL,.55,y+.05,.58,.58,palette[i]); tb(s,.55,y+.19,.58,.2,str(i+1),11,True,'FFFFFF',PP_ALIGN.CENTER); tb(s,1.45,y+.13,10.2,.36,x,15,False,INK)
    if i<len(items)-1: shape(s,MSO_SHAPE.DOWN_ARROW,.78,y+.66,.16,.25,'BFD4F7','BFD4F7')
  # Architecture and concept maps use a connected system view.
  elif 'architecture' in title.lower() or 'concept map' in title.lower():
   n=len(items); w=2.08
   for i,x in enumerate(items):
    x0=.65+i*2.48; shape(s,MSO_SHAPE.ROUNDED_RECTANGLE,x0,2.25,w,2.65,pale[i],palette[i]); tb(s,x0+.18,2.57,w-.36,.35,f'{i+1:02d}',13,True,palette[i]); tb(s,x0+.18,3.14,w-.36,1.18,x,15,True,INK,PP_ALIGN.CENTER)
    if i<n-1: shape(s,MSO_SHAPE.CHEVRON,x0+w+.08,3.23,.3,.55,'BFD4F7','BFD4F7')
  # General concept slides use a varied card grid.
  else:
   positions=[(.72,1.82,3.72,2.0),(4.8,1.82,3.72,2.0),(8.88,1.82,3.72,2.0),(2.76,4.18,3.72,2.0),(6.84,4.18,3.72,2.0)] if len(items)>=5 else [(.75,1.95,5.75,1.85),(6.75,1.95,5.75,1.85),(.75,4.15,5.75,1.85),(6.75,4.15,5.75,1.85)]
   for i,x in enumerate(items):
    px,py,pw,ph=positions[i]; shape(s,MSO_SHAPE.ROUNDED_RECTANGLE,px,py,pw,ph,'FFFFFF','D7E3F5'); shape(s,MSO_SHAPE.RECTANGLE,px,py,pw,.1,palette[i]); shape(s,MSO_SHAPE.OVAL,px+.2,py+.28,.46,.46,palette[i]); tb(s,px+.2,py+.39,.46,.18,str(i+1),10,True,'FFFFFF',PP_ALIGN.CENTER); tb(s,px+.82,py+.27,pw-1.05,1.15,x,14,True if len(x)<70 else False,INK)
 def step(title,kicker,n,text):
  s=shell(title,kicker); c=s.shapes.add_shape(MSO_SHAPE.OVAL,I(.75),I(2),I(1.15),I(1.15)); c.fill.solid(); c.fill.fore_color.rgb=C.from_string(BLUE); c.line.fill.background(); tb(s,.75,2.27,1.15,.5,str(n),24,True,'FFFFFF',PP_ALIGN.CENTER); tb(s,2.25,1.95,9.8,3.2,text,23,True if n==1 else False,INK)
 s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=C(255,255,255)
 if LOGO.exists(): s.shapes.add_picture(str(LOGO),I(.72),I(.42),I(.72),I(.72))
 tb(s,1.58,.62,4,.3,'TERTIARY INFOTECH ACADEMY',12,True,BLUE); tb(s,.75,1.5,5.75,2.0,TITLE,31,True); tb(s,.75,3.75,5.2,.5,f'Course Code: {CODE}',18,True,TEAL); tb(s,.75,4.48,5.6,.9,'Agentic AI Loop Engineering\nBuild · Verify · Ship',21,False,GREY); tb(s,.75,6.5,5.7,.3,'2 days · 15 hours · Intermediate · Version 2.0',12,False,GREY)
 if (VISUALS/'agentic-loop.png').exists(): s.shapes.add_picture(str(VISUALS/'agentic-loop.png'),I(6.45),I(1.12),I(6.55),I(3.68))
 bullets('How adults learn in this course','LEARNING DESIGN',['See a complete demonstration with trainer think-aloud','Practise the same decision with scaffolding','Solve a changed boundary independently','Explain evidence to a peer','Commit a transferable checklist or checkpoint'])
 bullets('The Agentic AI Engineering Loop','CORE MODEL',['Frame — define observable success','Plan — expose assumptions and file scope','Generate — authorise a small increment','Inspect + Verify — apply judgement and evidence','Correct + Commit — diagnose, fix and preserve'])
 bullets('Course outcomes','COURSE OVERVIEW',['Engineer durable goals, context and prompts','Build an accessible typed full-stack vertical slice','Generate tests, production builds and CI/CD gates','Create non-WSQ skills, hooks and bounded subagents','Implement MCP tools and orchestrate a release'])
 bullets('TaskFlow capstone architecture','COURSE OVERVIEW',['React + TypeScript user interface','Express API with schema validation','SQLite migrations and parameterised persistence','Shared contracts and automated tests','Docker, GitHub Actions and cloud release'])
 visual_story('See the whole system before you build','CAPSTONE VISUAL','full-stack-architecture.png',['Trace data end to end','Locate trust boundaries','Name the evidence at each layer'])
 topic_visuals={'01':('agentic-loop.png',['Human judgement directs the loop','Every stage leaves evidence','Commit only verified work']),'02':('full-stack-architecture.png',['Contracts connect the layers','Tests observe boundaries','Security travels with data']),'03':('cicd-pipeline.png',['Gates stop defects early','Deployments stay observable','Rollback remains deliberate']),'04':('multi-agent-mcp.png',['Delegation stays bounded','Tools enforce safe access','One lead reconciles the release'])}
 for tn,topic,concepts,labs in topics:
  s=shell(f'Topic {tn}',topic); tb(s,.7,1.9,2,1.3,tn,70,True,'DCE9FF'); tb(s,2.55,2.05,9.7,2.0,topic,34,True); tb(s,2.6,4.1,9,.8,f'{len(labs)} connected labs · concept → demonstration → practice → transfer',18,False,GREY)
  img,takes=topic_visuals[tn]; visual_story(topic,f'TOPIC {tn} · VISUAL STORY',img,takes)
  bullets('Why this topic matters',f'TOPIC {tn}',[f'{c}: connect knowledge to an engineering decision.' for c in concepts])
  bullets('Concept map',f'TOPIC {tn}',concepts)
  bullets('Common failure patterns',f'TOPIC {tn}',['A polished solution to the wrong problem','Chat-only context that disappears next session','Broad generated diffs that cannot be reviewed','Happy-path tests that miss real boundaries','Automation without ownership or rollback'])
  for labno,name,goal,target,practice in labs:
   bullets(f'Lab {labno} — {name}',f'TOPIC {tn} · LAB OVERVIEW',[goal,f'Primary target: {target}',practice,'Evidence: prompt + reviewed diff + command output','Independent challenge: change one boundary condition'])
   prompt=['GOAL — '+goal,'CONTEXT — relevant files and last verified checkpoint','CONSTRAINTS — scope, types, accessibility, security','DELIVERABLES — '+practice,'VERIFICATION — exact commands and expected behaviour']
   bullets('Reference prompt contract',f'LAB {labno}',prompt)
   procedures=lab_steps(labno,name,goal,target,practice)
   for start,end in [(0,4),(4,7),(7,10)]:
    bullets(f'Procedure · Steps {start+1}–{end}',f'LAB {labno}',[f'{i+1}. {procedures[i][0]} — {procedures[i][1]}' for i in range(start,end)])
   bullets('Expected result and evidence',f'LAB {labno}',[practice,'Save the final prompt and reviewed diff','Capture verification command and exit status','Record the boundary challenge result','Write one corrected assumption'])
   bullets('Troubleshoot, challenge and reflect',f'LAB {labno}',['Reproduce and minimise before prompting again','Name two hypotheses; run one discriminating test','Predict a changed boundary, then compare the result','Identify the highest-risk corrected assumption','Commit only after the full verification gate'])
 bullets('Capstone release review','COURSE CLOSE',['All 24 lab checkpoints are traceable','Lint, typecheck, unit, integration and build gates pass','MCP tools remain typed, read-only and project-contained','No secrets or prohibited programme content','Release note includes limitations and rollback'])
 bullets('Your repeatable practice','COURSE CLOSE',['Frame intent so success is observable','Plan before code and keep diffs reviewable','Inspect with human judgement','Verify with independent evidence','Correct deliberately and commit reversibly'])
 prs.save(OUT/f'{CODE} {TITLE}.pptx')

# Reference-standard deck: mirrors the supplied n8n deck's five-slide lab rhythm.
def build_ppt():
 prs=Presentation(); prs.slide_width=I(13.333); prs.slide_height=I(7.5)
 colors=[BLUE,TEAL,VIOLET,'F59E0B','E25555']; pales=['EAF2FF','ECFDF5','F3EEFF','FFF8E8','FFF4F4']
 def rect(s,x,y,w,h,fill='FFFFFF',line='DCE5F0',radius=True):
  q=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,I(x),I(y),I(w),I(h)); q.fill.solid(); q.fill.fore_color.rgb=C.from_string(fill); q.line.color.rgb=C.from_string(line); return q
 def circle(s,x,y,d,fill,text,size=12):
  q=s.shapes.add_shape(MSO_SHAPE.OVAL,I(x),I(y),I(d),I(d)); q.fill.solid(); q.fill.fore_color.rgb=C.from_string(fill); q.line.fill.background(); tb(s,x,y+d*.28,d,d*.35,text,size,True,'FFFFFF',PP_ALIGN.CENTER)
 def base(kicker,title):
  s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=C(255,255,255)
  tb(s,.68,.3,11.7,.22,kicker.upper(),8,True,TEAL); tb(s,.68,.68,11.9,.82,title,18,True,INK)
  s.shapes.add_shape(MSO_SHAPE.RECTANGLE,I(.68),I(6.94),I(11.95),I(.012)).fill.solid(); s.shapes[-1].fill.fore_color.rgb=C.from_string('DCE5F0'); s.shapes[-1].line.fill.background()
  tb(s,.68,7.05,5.8,.14,f'{CODE} - {TITLE}',7,False,GREY); tb(s,5.3,7.05,4.8,.14,'Copyright 2026 Tertiary Infotech Academy Pte Ltd',7,False,GREY,PP_ALIGN.CENTER); tb(s,12.15,7.02,.4,.16,str(len(prs.slides)),7,False,GREY,PP_ALIGN.RIGHT); return s
 def card(s,x,y,w,h,color,label,body,num=None):
  rect(s,x,y,w,h,'FFFFFF','DCE5F0'); rect(s,x,y,.08,h,color,color,False)
  if num is not None: circle(s,x+.23,y+.24,.42,color,str(num),9); ox=x+.78
  else: ox=x+.3
  tb(s,ox,y+.22,w-(ox-x)-.22,.26,label.upper(),9,True,color); tb(s,x+.3,y+.72,w-.6,h-.88,body,11,False,INK)
 def loop_row(s,y=5.45):
  names=['Define','Build','Observe','Evaluate','Improve']; xs=[.85,2.25,3.65,5.05,6.45]
  for i,(x,n) in enumerate(zip(xs,names)):
   circle(s,x,y,.38,TEAL if i==4 else BLUE,str(i+1),8); tb(s,x-.2,y+.48,.78,.18,n,7,False,GREY,PP_ALIGN.CENTER)
   if i<4:
    ln=s.shapes.add_shape(MSO_SHAPE.CHEVRON,I(x+.45),I(y+.11),I(.8),I(.15)); ln.fill.solid(); ln.fill.fore_color.rgb=C.from_string('DCE5F0'); ln.line.fill.background()
 def target_slide(tn,labno,name,goal,target,practice):
  s=base('Lab Target',f'Lab {labno}: {name}')
  rect(s,.72,1.55,7.75,4.45,'F5F8FC','DCE5F0'); circle(s,1.12,2.0,1.0,BLUE,labno,15)
  tb(s,2.38,1.82,5.55,.25,'YOU WILL BUILD',9,True,BLUE); tb(s,2.38,2.32,5.45,1.35,goal,18,True,INK)
  tb(s,2.38,4.0,5.2,.22,'DELIVERABLE',8,True,TEAL); tb(s,2.38,4.38,5.4,.8,practice,11,False,GREY)
  entries=[('TIMEBOX','25-40 minutes guided practice','F59E0B'),('MODE','Plan → generate → inspect → verify',VIOLET),('EVIDENCE',f'Prompt, diff and checkpoint in {target}',TEAL)]
  for i,(a,b,c) in enumerate(entries): card(s,8.8,1.66+i*1.42,3.85,1.18,c,a,b)
  tb(s,.78,6.3,7.8,.25,'Success is demonstrated through visible evidence and learner explanation.',10,True,VIOLET)
 def ideas_slide(labno,name,goal,target,practice):
  s=base('Concept Map',f'Lab {labno}: Four Ideas That Make It Work')
  ideas=[('Outcome boundary',goal),('Primary target',target),('Engineering practice',practice),('Independent evidence','A changed boundary condition plus a repeatable verification result.')]
  pos=[(.72,1.55),(6.75,1.55),(.72,3.78),(6.75,3.78)]
  for i,((label,body),(x,y)) in enumerate(zip(ideas,pos)): card(s,x,y,5.85,1.85,colors[i],label,body,i+1)
  tb(s,.78,6.18,9.2,.22,'Connect the four ideas before touching the implementation.',9,False,GREY)
 def sequence_slide(labno,name,goal,target,practice):
  s=base('Workflow Plan',f'Lab {labno}: Build Sequence'); steps=lab_steps(labno,name,goal,target,practice)
  selected=[steps[0],steps[2],steps[4],steps[5],steps[6],steps[9]]; pos=[(.72,1.52),(6.75,1.52),(.72,2.86),(6.75,2.86),(.72,4.2),(6.75,4.2)]
  for i,((phase,detail),(x,y)) in enumerate(zip(selected,pos)):
   rect(s,x,y,5.85,1.02,'FFFFFF','DCE5F0'); circle(s,x+.18,y+.22,.42,BLUE if i%2==0 else TEAL,str(i+1),8); tb(s,x+.78,y+.18,4.8,.22,phase,9,True,INK); tb(s,x+.78,y+.5,4.75,.35,detail,8,False,GREY)
  tb(s,.78,5.74,8.8,.2,'Build one observable behaviour at a time; inspect before adding the next.',9,True,VIOLET)
 def gate_slide(labno,name,goal,target,practice):
  s=base('Quality Gate',f'Lab {labno}: Prove, Evaluate, Improve')
  rect(s,.72,1.5,7.15,4.9,'F5F8FC','DCE5F0')
  checks=[practice,'The diff stays within the declared target.','Lint, typecheck and focused tests pass.','One edge case is run and the correction is recorded.']
  for i,x in enumerate(checks): circle(s,1.05,1.88+i*1.02,.38,TEAL,'✓',10); tb(s,1.65,1.84+i*1.02,5.75,.5,x,11,False,INK)
  rect(s,8.2,1.5,4.4,4.9,'FFFFFF','DCE5F0'); tb(s,8.55,1.82,2.2,.2,'EVIDENCE STACK',8,True,BLUE)
  for i,(a,b) in enumerate([('Input','Saved prompt'),('Trace','Reviewed diff'),('Output','Command result'),('Decision','Corrected assumption')]):
   circle(s,8.55,2.35+i*.8,.34,VIOLET,str(i+1),7); tb(s,9.08,2.29+i*.8,1.05,.18,a,8,True,VIOLET); tb(s,10.25,2.29+i*.8,1.85,.25,b,8,False,GREY)
  tb(s,8.55,5.63,.55,.18,'PASS',7,True,TEAL); rect(s,9.25,5.67,2.55,.12,TEAL,TEAL)
 def troubleshoot_slide(labno,name,goal,target,practice):
  s=base('Troubleshooting',f'Lab {labno}: Diagnose Before Rebuilding')
  headers=[('SYMPTOM','E25555'),('LIKELY CAUSE','F59E0B'),('TARGETED FIX',TEAL)]
  for i,(h,c) in enumerate(headers): rect(s,.72+i*4.02,1.5,3.83,.45,c,c); tb(s,.72+i*4.02,1.64,3.83,.15,h,8,True,'FFFFFF',PP_ALIGN.CENTER)
  rows=[('Agent changes too much','Scope is broader than the named target','Stop, restore and authorise named files only.'),('Verification is inconsistent','Commands or data are not deterministic','Reset state and rerun the same narrow check.'),('Output looks correct but fails','The test misses the real system boundary','Exercise the actual boundary and add a regression test.')]
  for r,row in enumerate(rows):
   for c,txt in enumerate(row):
    fill=['FFF4F4','FFF8E8','ECFDF5'][c]; rect(s,.72+c*4.02,2.1+r*1.14,3.83,.88,fill,'DCE5F0'); tb(s,.95+c*4.02,2.34+r*1.14,3.35,.38,txt,9,False,INK)
  tb(s,.78,5.82,9.4,.22,'Change one variable, rerun the same test, and compare the trace.',9,True,VIOLET)
 # Cover
 s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=C(255,255,255)
 tb(s,.68,.36,5,.2,'TERTIARY INFOTECH ACADEMY',8,True,TEAL); tb(s,.68,1.15,6.1,1.35,TITLE,27,True,INK); tb(s,.68,3.03,5.8,.55,'Agentic AI loop engineering for adult full-stack training',14,False,GREY); tb(s,.68,5.55,5.8,.25,f'{CODE}  |  Version 3.0  |  UEN 201200696W',9,False,GREY)
 rect(s,7.25,.7,5.35,5.75,'F5F8FC','DCE5F0'); nodes=[('FRONTEND',8.0,1.35,BLUE),('API',10.15,1.35,VIOLET),('DATA',8.0,3.2,TEAL),('DELIVERY',10.15,3.2,'F59E0B'),('REVIEW',9.08,4.9,'E25555')]
 for label,x,y,c in nodes: circle(s,x,y,.82,c,''); tb(s,x-.15,y+.95,1.12,.2,label,8,True,c,PP_ALIGN.CENTER)
 # Outcomes + loop
 s=base('Course Outcomes','What You Will Be Able to Build')
 outs=[('AGENT CONTEXT','Goals, constraints and durable instructions'),('FULL STACK','React, Express, SQLite and shared types'),('QUALITY','Tests, accessibility and safe errors'),('DELIVERY','CI/CD, containers and rollback'),('ORCHESTRATION','Skills, subagents and MCP tools')]
 for i,(label,body) in enumerate(outs): card(s,.72+(i%3)*4.02,1.55+(i//3)*2.0,3.82,1.55,colors[i],label,body,i+1)
 s=base('Engineering Method','The Agentic AI Loop'); loop_row(s,2.3); tb(s,.82,3.35,11.1,.75,'Define the outcome. Build a small increment. Observe the trace. Evaluate against evidence. Improve deliberately.',17,True,INK,PP_ALIGN.CENTER); rect(s,2.1,4.55,9.0,.85,'F5F8FC','DCE5F0'); tb(s,2.45,4.83,8.3,.25,'Every revision must leave the repository easier to understand and safer to change.',11,False,GREY,PP_ALIGN.CENTER)
 for tn,topic,concepts,labs in topics:
  s=base(f'Topic {tn}',f'Topic {tn} - {topic}'); rect(s,.72,1.55,8.0,2.55,'F5F8FC','DCE5F0'); tb(s,1.0,1.9,7.4,1.05,'; '.join(concepts)+'.',13,False,INK); tb(s,11.0,1.45,1.2,.7,tn,42,True,'E8EEF6',PP_ALIGN.CENTER); loop_row(s,4.55); rect(s,9.45,4.25,3.1,1.75,'FFFFFF','DCE5F0'); tb(s,9.77,4.5,1.2,.18,'LEARNING ARC',7,True,VIOLET); tb(s,9.77,4.87,2.2,.8,'Concept\nPractice\nEvidence\nReflection',10,True,INK)
  for labno,name,goal,target,practice in labs:
   target_slide(tn,labno,name,goal,target,practice); ideas_slide(labno,name,goal,target,practice); sequence_slide(labno,name,goal,target,practice); gate_slide(labno,name,goal,target,practice); troubleshoot_slide(labno,name,goal,target,practice)
 s=base('Course Close','Build with Evidence. Operate with Confidence.'); rect(s,1.05,1.7,11.2,3.7,'F5F8FC','DCE5F0'); tb(s,1.55,2.2,10.2,.5,'A polished output is only the beginning.',20,True,INK,PP_ALIGN.CENTER); tb(s,1.65,3.1,10.0,.65,'Professional full-stack systems are observable, testable, guarded, recoverable, and documented.',14,False,GREY,PP_ALIGN.CENTER)
 for i,x in enumerate(['TRACE','TEST','REVIEW','RECOVER','DOCUMENT']): rect(s,2.0+i*1.85,4.35,1.55,.42,colors[i],colors[i]); tb(s,2.0+i*1.85,4.48,1.55,.15,x,7,True,'FFFFFF',PP_ALIGN.CENTER)
 tb(s,5.0,5.95,3.3,.2,'www.tertiarycourses.com.sg',8,True,BLUE,PP_ALIGN.CENTER)
 prs.save(OUT/f'{CODE} {TITLE}.pptx')

def build_lp():
 d=base_doc('Lesson Plan — Agentic AI Loop Engineering')
 d.add_heading('Course Information',0); para(d,'Intermediate | 2 days | 15 instructional hours | Instructor-led demonstration, guided practice and independent challenge')
 d.add_heading('Adult Learning Design',0); para(d,'Each block uses activation of prior experience, concise concept input, trainer think-aloud, guided lab practice, independent boundary challenge, peer explanation and reflective checkpoint. The trainer adapts scaffolding based on observable evidence rather than pace alone.')
 d.add_heading('Learning Outcomes',0)
 for x in ['Engineer durable agent goals, context and repository instructions.','Implement and test an accessible typed full-stack application.','Automate production builds, CI/CD gates and rollback evidence.','Create namespaced skills, hooks, subagents and secure MCP tools.']: d.add_paragraph(x,style='List Bullet')
 schedules=[('Day 1 — Context, architecture and implementation',[
 ('09:30–10:00','Orientation, prior-experience activation and loop demonstration',30),('10:00–11:15','Topic 1 concepts + Labs 1.1–1.3',75),('11:15–13:00','Labs 1.4–1.6: brief, instructions and scaffold',105),('13:00–13:30','Lunch',30),('13:30–15:00','Topic 2 concepts + Labs 2.1–2.4',90),('15:00–15:15','Tea break',15),('15:15–17:00','Labs 2.5–2.9: contracts, API, data and integration',105),('17:00–17:30','Labs 2.10–2.11, review and checkpoint',30)]),
 ('Day 2 — Delivery, reusable agents and orchestration',[
 ('09:30–10:00','Retrieval practice and checkpoint recovery',30),('10:00–11:30','Topic 3 concepts + Labs 3.1–3.3',90),('11:30–13:00','Labs 3.4–3.5: CI and failure injection',90),('13:00–13:30','Lunch',30),('13:30–14:30','Labs 3.6–3.7: deploy, observe and roll back',60),('14:30–15:00','Topic 4 concepts and trainer MCP demonstration',30),('15:00–15:15','Tea break',15),('15:15–16:45','Labs 4.1–4.5: skills, hooks, subagents and MCP',90),('16:45–17:30','Lab 4.6 capstone release, peer explanation and close',45)])]
 d.add_heading('Detailed Schedule',0)
 for day,rows in schedules:
  d.add_heading(day,1); t=d.add_table(rows=1,cols=3); t.style='Table Grid'; t.alignment=WD_TABLE_ALIGNMENT.CENTER
  for i,x in enumerate(['Time','Learning Activity','Minutes']): t.cell(0,i).text=x; shade(t.cell(0,i),BLUE)
  total=0
  for tm,act,m in rows:
   cells=t.add_row().cells; total+=m
   for i,x in enumerate([tm,act,str(m)]): cells[i].text=x
   if act=='Lunch' or 'break' in act.lower():
    for c in cells: shade(c,'FFF4CC')
  assert total==480
  para(d,'Instructional/contact time excluding lunch: 450 minutes (7.5 hours).')
 d.add_heading('Facilitation and Learning Checks',0); para(d,'Use observation checklists, prompt/diff walkthroughs, command output, peer explanation, failure injection and capstone demonstration. Provide extra scaffolding through partially completed prompt contracts; extend faster learners with changed boundaries and security review.')
 d.save(OUT/f'Lesson Plan - {CODE} {TITLE}.docx')

if __name__=='__main__':
 OUT.mkdir(exist_ok=True); md=build_labs_and_md(); build_guide(md); build_ppt(); build_lp(); print('v2 generated:',sum(1 for _ in all_labs()),'labs')
