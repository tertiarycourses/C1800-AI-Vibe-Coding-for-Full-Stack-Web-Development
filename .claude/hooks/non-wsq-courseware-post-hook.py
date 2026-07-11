#!/usr/bin/env python3
from pathlib import Path
import re, sys
from pptx import Presentation
root=Path(sys.argv[1] if len(sys.argv)>1 else '.').resolve()
deliverables=[root/'LEARNER-GUIDE.md',*sorted((root/'labs').rglob('*.md'))]
prohibited=re.compile(r'\b(SSG|TRAQOM)\b|SkillsFuture funding',re.I); errors=[]
for p in deliverables:
    if not p.exists(): errors.append(f'missing: {p}'); continue
    for n,line in enumerate(p.read_text(errors='ignore').splitlines(),1):
        if prohibited.search(line): errors.append(f'{p}:{n}: prohibited programme language')
for folder in [root/'.codex/skills',root/'.claude/commands',root/'.claude/hooks',root/'.claude/agents']:
    if folder.exists():
        for p in folder.iterdir():
            if not p.name.startswith('non-wsq-'): errors.append(f'{p}: custom tooling lacks non-wsq- prefix')
required=[root/'courseware'/f for f in ['C1800 AI Vibe Coding for Full Stack Web Development.pptx','C1800 AI Vibe Coding for Full Stack Web Development - Learner Guide.docx','Lesson Plan - C1800 AI Vibe Coding for Full Stack Web Development.docx']]
for p in required:
    if not p.exists() or p.stat().st_size<10000: errors.append(f'missing or undersized: {p}')
lab_count=len(list((root/'labs').rglob('*.md')))
if lab_count < 20: errors.append(f'only {lab_count} labs; at least 20 required')
if required[0].exists() and len(Presentation(required[0]).slides) <= 100: errors.append('PPT must exceed 100 slides')
if errors: print('\n'.join(errors)); raise SystemExit(1)
print(f'non-wsq post-hook passed: {lab_count} labs, {len(Presentation(required[0]).slides)} slides and 3 primary artifacts')
